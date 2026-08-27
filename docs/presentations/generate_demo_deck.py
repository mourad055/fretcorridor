#!/usr/bin/env python3
"""Génère la présentation PowerPoint FretCorridor — Parcours démo (Flysoft Engineering)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
LOGO = ROOT / "assets" / "flysoft-engineering-logo.png"
OUTPUT = ROOT / "FretCorridor-Parcours-Demo-Flysoft.pptx"

# Palette Flysoft (logo + charte sobre)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)
BLUE = RGBColor(0x1E, 0x5A, 0xA8)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x64, 0x64, 0x64)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0x8A, 0x4B)
RED = RGBColor(0xC6, 0x28, 0x28)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_brand_bar(slide) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_logo(slide) -> None:
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.55), Inches(0.28), width=Inches(2.35))


def add_footer(slide, text: str = "Flysoft Engineering · FretCorridor · Confidentiel démo") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.35))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY


def add_title_block(slide, title: str, subtitle: str | None = None) -> None:
    add_brand_bar(slide)
    add_logo(slide)
    tbox = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(9.5), Inches(1.2))
    tf = tbox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(9.5), Inches(0.8))
        stf = sbox.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(16)
        sp.font.color.rgb = GRAY


def add_bullets(slide, items: list[str], top: float = 2.0, left: float = 0.75, width: float = 11.5) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_step_card(slide, x: float, y: float, num: int, title: str, subtitle: str, color: RGBColor) -> None:
    w, h = 1.85, 1.55
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2.5)

    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.12), Inches(y + 0.12), Inches(0.42), Inches(0.42))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.text = str(num)
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cp.font.size = Pt(14)
    cp.font.bold = True
    cp.font.color.rgb = WHITE

    tbox = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.58), Inches(w - 0.16), Inches(0.55))
    ttf = tbox.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(11)
    tp.font.bold = True
    tp.font.color.rgb = DARK
    tp.alignment = PP_ALIGN.CENTER

    sbox = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 1.05), Inches(w - 0.16), Inches(0.45))
    stf = sbox.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(9)
    sp.font.color.rgb = GRAY
    sp.alignment = PP_ALIGN.CENTER


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_brand_bar(slide)
    add_logo(slide)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.4), SLIDE_W, Inches(2.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(11), Inches(1.2))
    tf = tbox.text_frame
    tf.text = "FretCorridor"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.65), Inches(11), Inches(0.9))
    stf = sbox.text_frame
    stf.text = "Parcours de démonstration — Plateforme logistique corridor"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(22)
    sp.font.color.rgb = WHITE

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8), Inches(0.5))
    mtf = meta.text_frame
    mtf.text = "Flysoft Engineering  ·  Août 2026  ·  Web + Mobile + Moteur"
    mp = mtf.paragraphs[0]
    mp.font.size = Pt(14)
    mp.font.color.rgb = GRAY


def slide_content(prs: Presentation, title: str, subtitle: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_title_block(slide, title, subtitle)
    add_bullets(slide, bullets)
    add_footer(slide)


def slide_journey(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_title_block(slide, "Parcours bout-en-bout", "6 étapes de la démo — de l'inscription au reversement transporteur")

    steps = [
        (0, "Connexion\n& KYC", "Onboarding", GRAY),
        (1, "Demande\nchargeur", "Marketplace", RED),
        (2, "Capacité\ntransporteur", "Offre", RED),
        (3, "Appariement\nMoteur", "Matching", ORANGE),
        (4, "Exécution\nmission", "GPS · ETA", BLUE),
        (5, "Paiement\n& reversement", "Finance", GREEN),
    ]
    x0, y = 0.45, 2.15
    gap = 2.05
    for i, (num, title, sub, color) in enumerate(steps):
        add_step_card(slide, x0 + i * gap, y, num, title.replace("\n", " "), sub, color)
        if i < len(steps) - 1:
            ax = x0 + i * gap + 1.92
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(ax), Inches(y + 0.62), Inches(0.18), Inches(0.22)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

    legend = slide.shapes.add_textbox(Inches(0.65), Inches(4.2), Inches(12), Inches(2.2))
    ltf = legend.text_frame
    ltf.word_wrap = True
    ltf.text = (
        "Architecture réelle : 3 applications (Client, Chauffeur/Transporteur, Portail Web) "
        "+ Gateway + microservices Java + Kafka + PostgreSQL + Valhalla (ETA).\n"
        "Le Moteur (OPT/MAT/GEO/TRK) travaille en arrière-plan — invisible pour l'utilisateur final."
    )
    lp = ltf.paragraphs[0]
    lp.font.size = Pt(14)
    lp.font.color.rgb = DARK
    add_footer(slide)


def slide_actors_table(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_title_block(slide, "Rôle de chaque acteur", "Qui fait quoi pendant la démonstration")

    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(2.0), Inches(12.2), Inches(4.5))
    table = table_shape.table

    headers = ["Acteur / interface", "Rôle métier", "Étapes actives", "Ce qu'on montre"]
    data = [
        [
            "App Client (mobile)",
            "Chargeur — exprime un besoin de transport",
            "0 → 1 → 4",
            "Inscription · KYC · publier demande · propositions · suivi",
        ],
        [
            "App Chauffeur / Transporteur",
            "Offre & exécution terrain",
            "0 → 2 → 3 → 4 → 5",
            "Véhicule · capacité · propositions · mission · solde",
        ],
        [
            "Portail Web Bureau / Admin",
            "Supervision du corridor",
            "4 → 5",
            "Carte GPS · observatoire · missions · notifs · rapport financier",
        ],
        [
            "Moteur FretCorridor (backend)",
            "Appariement intelligent invisible",
            "1 → 2 → 3 → 4",
            "Matching L0/L1 · coût 7 termes · ETA · événements Kafka",
        ],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK

    add_footer(slide)


def slide_demo_script(prs: Presentation, phase: str, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_title_block(slide, title, phase)
    add_bullets(slide, bullets, top=2.1)
    add_footer(slide)


def build() -> Path:
    prs = Presentation()
    set_slide_size(prs)

    slide_title(prs)

    slide_content(
        prs,
        "Objectif de la démo",
        "Ce que les interlocuteurs doivent comprendre",
        [
            "FretCorridor est une plateforme logistique corridor réelle — pas une maquette statique.",
            "Trois interfaces utilisateur communiquent avec un moteur d'appariement intelligent.",
            "Le chargeur publie un besoin, le transporteur déclare une offre, le système calcule le meilleur match.",
            "La mission s'exécute avec suivi GPS et ETA, puis le paiement est réconcilié côté transporteur et bureau.",
            "Phrase clé : « Offre ↔ Demande ↔ Mission ↔ Paiement » sur infrastructure microservices.",
        ],
    )

    slide_journey(prs)
    slide_actors_table(prs)

    slide_content(
        prs,
        "Étapes 0 à 2 — Préparation & marketplace",
        "Onboarding, demande et offre",
        [
            "Étape 0 — Connexion & KYC : inscription légère puis profil niveau 1 (identité + pièce justificative).",
            "Étape 1 — Chargeur publie une demande (App Client) : axe Douala → Yaoundé, emballage, fenêtre horaire.",
            "Service : service-mkt → événement Kafka demande-publiee → file d'attente service-opt.",
            "Étape 2 — Transporteur déclare une capacité (App Chauffeur) : véhicule enregistré, même axe, GPS activé.",
            "Service : service-cap → Kafka capacite-declaree → matching possible dès qu'une paire existe sur l'axe.",
        ],
    )

    slide_content(
        prs,
        "Étapes 3 à 5 — Moteur, exécution & finance",
        "Matching, mission et reversement",
        [
            "Étape 3 — Appariement Moteur (invisible) : cycle ~15 s, filtrage L0 H3, coût 7 termes, Kuhn-Munkres.",
            "Services : service-opt · service-mat · service-geo · Valhalla → proposition émise vers service-mkt.",
            "Étape 4 — Exécution : acceptation proposition → mission → positions GPS → ETA → carte Bureau.",
            "Services : service-exe · service-flt · service-trk · service-bur (Kafka position-eta).",
            "Étape 5 — Paiement : solde transporteur, rapport financier Bureau, notifications IN_APP.",
        ],
    )

    slide_content(
        prs,
        "Architecture technique",
        "3 apps · 1 gateway · microservices",
        [
            "App Client → appels directs : service-ida (8081), service-mkt (8089), service-geo (8084)…",
            "App Chauffeur/Transporteur → Gateway (8082) : capacités, véhicules, missions, notifications.",
            "Portail Web (8099) → Gateway → services métier (Bureau, Admin, Transporteur web).",
            "Moteur : OPT + MAT + GEO + TRK tournent en conteneurs Docker (23 services au total).",
            "Données : PostgreSQL · Kafka · MinIO (pièces KYC) · Valhalla (routage ETA).",
        ],
    )

    slide_demo_script(
        prs,
        "Phase A — Recommandée en ouverture (données seedées, fiable)",
        "Script présentateur · ~10 minutes",
        [
            "1. Portail Bureau (+237600000001 / 1234) — Observatoire, carte positions, 2 notifications IN_APP.",
            "2. Transporteur web/mobile (+237696000001 / 1234) — 11 missions, capacités, solde 257 000 FCFA.",
            "3. Admin (+237600000003 / 1234) — KYC (ouvrir pièce, télécharger), dossiers, tenants.",
            "4. Message : « L'écosystème est branché en réel — ce que vous voyez vient des APIs et de la base. »",
        ],
    )

    slide_demo_script(
        prs,
        "Phase B — Scénario live (optionnel)",
        "Nouvelles inscriptions · ~10 minutes · prévoir 60–90 s pour le matching",
        [
            "1. Nouveau chargeur (App Client) + nouveau transporteur (App Chauffeur) — numéros neufs.",
            "2. KYC niveau 1 des deux côtés (identité + photo pièce).",
            "3. Transporteur : ajouter véhicule → déclarer capacité axe Douala → Yaoundé.",
            "4. Chargeur : publier demande sur le même axe (utiliser le sélecteur d'axe, pas la saisie libre).",
            "5. Attendre 60–90 s → proposition côté transporteur → acceptation → mission visible Bureau.",
        ],
    )

    slide_content(
        prs,
        "Accès & comptes démo",
        "URLs et identifiants",
        [
            "Portail Web : http://localhost:8099  ·  Gateway API : http://localhost:8082",
            "Bureau : +237600000001 — PIN 1234",
            "Transporteur : +237696000001 — PIN 1234",
            "Admin : +237600000003 — PIN 1234",
            "Mobile physique : lancer ./scripts/run_dev.sh (IP WiFi, pas localhost).",
            "Stack Docker : 23 services — attendre ~2 min après démarrage avant matching live.",
        ],
    )

    slide_content(
        prs,
        "Messages clés & plan B",
        "Anticiper les questions du manager",
        [
            "Ce n'est pas instantané : le matching tourne par cycles planifiés (~15 secondes).",
            "KYC niveau 1 obligatoire avant toute publication de demande.",
            "Véhicule obligatoire avant déclaration de capacité transporteur.",
            "Plan B : si le live tarde, revenir aux données seedées (Phase A) — toujours impressionnant.",
            "Après la démo : compléter le CDC module par module sans régression sur ce parcours.",
        ],
    )

    # Slide clôture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLUE
    add_brand_bar(slide)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.9), Inches(1.2), width=Inches(3.5))
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.2))
    tf = tbox.text_frame
    tf.text = "Merci — Questions ?"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.8))
    stf = sbox.text_frame
    stf.text = "Flysoft Engineering · FretCorridor · Plateforme logistique intelligente"
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(16)
    sp.font.color.rgb = WHITE

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Présentation générée : {path}")
